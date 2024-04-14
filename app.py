import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import streamlit as st
from openai import OpenAI
from pptx.dml.color import RGBColor
from moviepy.editor import concatenate_videoclips, VideoFileClip, ImageSequenceClip, AudioFileClip
from pdf2image import convert_from_path

def get_text(lesson_name, grade):
    from googlesearch import search

    def search_and_filter(term):
        try:
            results = search(term, num=10, stop=10, pause=2)
            filtered_links = [link for link in results if link.startswith('https://ncert.nic.in')]
            return filtered_links
        except Exception as e:
            print("An error occurred:", str(e))
            return []

    search_term = "NCERT " + lesson_name + " " + grade + " filetype:pdf"
    print(search_term)
    filtered_links = search_and_filter(search_term)

    if filtered_links:
        import requests

        def download_pdf(url, filename):
            response = requests.get(url)
            with open(filename, 'wb') as f:
                f.write(response.content)
            print(f"PDF downloaded successfully as {filename}")

        url = f"{filtered_links[0]}"
        print(url)
        filename = "session" + ".pdf"
        download_pdf(url, filename)
    else:
        print("No links found.")
        st.warn("Enter a valid NCERT lesson name!")

    from langchain.text_splitter import CharacterTextSplitter
    from langchain_community.document_loaders import PyPDFLoader

    book_name = "session.pdf"
    loader = PyPDFLoader(f"{book_name}")
    documents = loader.load()
    text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=0)
    docs = text_splitter.split_documents(documents)
    with open("text.txt", "w") as f:
        for i in range(0, len(docs)):
            f.write(docs[i].page_content.replace("\t", " "))

    with open("text.txt", "r") as f:
        text = f.read()
    return text

def generate_subheadings(lesson_name,grade, api_key):
    lesson_content = get_text(lesson_name, grade)   
    book_name = "NCERT"
    grade = grade

    system_content = f"""
    You are an AI assistant specialized in preparing teaching material for students of all ages and grades. You will be given an entire lesson and you have to figure out the
    important subheadings and topics that will be useful for preparing material for the students. The subheadings and subtopics must be only from the lesson content you're provided. Don't make up your own.
    lesson name: {lesson_name}
    book name: {book_name}
    grade: {grade}
    give them in a number format like this:
    1. <subheading>
    1.1 <subheading>
    2. <subheading>
    """

    user_content = f""" This is the lesson content: {lesson_content} """

    from openai import OpenAI
    client = OpenAI(api_key = api_key)
    response = client.chat.completions.create(
    model="gpt-4-turbo",
    messages=[
    {"role": "system", "content": system_content},
    {"role": "user", "content": user_content},]
    )
    return response.choices[0].message.content


def app():

    api_key = st.text_input("Enter your API key", type="password")
    client = OpenAI(api_key=api_key)
    if os.path.exists("session_folder"):
        shutil.rmtree("session_folder")
        os.mkdir("session_folder")
    else:
        os.mkdir("session_folder")

    system_message = """
    You are a teacher with decades of experience in teaching students of all levels and languages. You are 
    talented and help students to learn and understand the subject matter. You will be given a 
    lesson name, the book name, the age group of the student. You are tasked with creating a PPT for the students based on the lesson content you receive.
    Each slide of the ppt must explain a portion of the lesson perfectly. 
    Follow these guidelines strictly:
    Short description must be the overview of what you're about to teach and it has to be atleast 5 - 10 lines long.
    Long description must be the detailed explanation of the topic you're teaching and it has to be atleast 30 - 50 lines long.
    You will be given a list of sub headings. You must definetly explain these sub headings in your ppt according to the lesson.

    Example: Slide <number>: <title> (this can be the topic you have chosen to teach.)
    short description: <short overview about the topic you're teaching, including the important keywords. This description has to be 5 - 10 lines"
    long description: <long description including the important keywords. This description has to be 30 - 50 lines or more. Explain precisely the topic you're explaining"
    
    DON'T LEAVE A LINE ANY WHERE IN YOUR RESPONSE.

    make sure to use vocabulary and language that is simple and very easy to understand.
    At the end of last slide, do not give any conclusion or summary. Just end the presentation. Thanks.
    """

    st.title("Ekai")
    st.markdown("""
                > High quality narration videos in minutes!
                """)

    class_options = ["Class " + str(i) for i in range(1, 13)]
    book_name = st.selectbox('Select book name', ["NCERT"])
    grade = st.selectbox('Select Class', class_options)
    lesson_name = st.text_input('Enter lesson name')

    age_group_options = ["6-7 years", "7-8 years", "8-9 years", "9-10 years", "10-11 years", "11-12 years",
                         "12-13 years", "13-14 years", "14-15 years", "15-16 years", "16-17 years", "17-18 years"]
    age_group = st.selectbox('Select Age Group', age_group_options)
    speed = st.select_slider('Select Speed', options=[0.75, 1, 1.25, 1.5], value=1)
    voice = st.selectbox('Select Voice', ("alloy", "echo", "fable", "nova", "onyx", "shimmer"))
    count = 0
    # slide_background_color = st.color_picker("Choose Slide Background Color", value="#D5E1DD")
    # fill_gradient_color_start = st.color_picker("Choose Fill Gradient Start Color", value="#9FDAC9")
    # fill_gradient_color_end = st.color_picker("Choose Fill Gradient End Color", value="#2E644E")


    if st.button("Let's go"):
        script = get_text(lesson_name, grade)
        sub_headings = generate_subheadings(lesson_name, grade, api_key)
        print(sub_headings)
        
        lines = [line.strip() for line in sub_headings.split("\n") if line.strip()]  # Remove empty lines
        subheading_string = "\n".join(lines) + "\n"  # Join non-empty lines
        print(subheading_string)

        sentences = [line.split(". ", 1)[1] if ". " in line else line for line in subheading_string.split("\n")]
        # chunks = [sentences[i:i + 4] for i in range(0, len(sentences), 4)]
        chunks = [sentences[i:i + 7] for i in range(0, len(sentences), 7)]

        for chunk in chunks:
            subheadings = ""
            for i in chunk:
                subheadings += i + "\n"

            user_message = f"""
            Lesson Name: {lesson_name}
            Book Name: {book_name}
            Age Group: {age_group}
            grade: {grade}
            Sub headings: {subheadings}
            script: {script}
            """

            if len(user_message) / 4 > 512000:
                st.error("Please keep the script under 512000 characters")
            else:
                user_message = f"""
                Lesson Name: {lesson_name}
                Book Name: {book_name}
                Age Group: {age_group}
                Sub headings: {subheadings}
                Lesson content: {script}
                """

                st.success("Processing")
                response = client.chat.completions.create(
                    model="gpt-4-0125-preview",
                    messages=[
                        {"role": "system", "content": f"{system_message}"},
                        {"role": "user", "content": f"{user_message}"},
                    ],
                    max_tokens=2000
                )

                full_script = response.choices[0].message.content[:-10]
                print(full_script)
                sections = full_script.split("\n\n")

                for k in range(0, len(sections)):
                    slide_title = sections[k].split("\n")[0]
                    short_description = sections[k].split("\n")[1][18:]
                    try:
                        long_description = sections[k].split("\n")[2][17:]
                    except:
                        long_description = short_description
                    speech_file_path = "session_folder/audio.mp3"

                    try:
                        os.remove("session_folder/presentation.pdf")
                    except:
                        pass

                    try:
                        os.remove("session_folder/presentation.pptx")
                    except:
                        pass

                    try:
                        os.remove("audio.mp3")
                    except:
                        print("failed to remove")

                    response = client.audio.speech.create(
                        model="tts-1",
                        voice=voice,
                        input=f"{long_description}",
                        speed=speed
                    )

                    response.stream_to_file(speech_file_path)
                       # Add options for slide background color
                    
                    prs = Presentation()
                    
                    slide_background = RGBColor(0xD5, 0xE1, 0xDD)  # Mint cream color
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = slide_background

                    title = slide.shapes.title
                    title.text = str(slide_title)[8:]
                    title.text_frame.paragraphs[0].font.bold = True  # Make the title bold
                    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # black color
                    title.text_frame.paragraphs[0].font.size = Pt(28)  # Increase font size

                    content = slide.placeholders[1]
                    content.text = short_description
                    content.text_frame.paragraphs[0].font.size = Pt(18)  # Set the font size to 18
                    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x2E, 0x64, 0x4E)  # Forest green color

                    fill = title.fill
                    fill.gradient()
                    fill.gradient_stops[0].color.rgb = RGBColor(0x9F, 0xDA, 0xC9)  # Light green
                    fill.gradient_stops[1].color.rgb = RGBColor(0x2E, 0x64, 0x4E)  # Dark green

                    content.shadow.inherit = False
                    content.shadow.visible = True
                    content.shadow.blur_radius = Pt(5)
                    content.shadow.offset_x = Inches(0.1)
                    content.shadow.offset_y = Inches(0.1)



                    presentation_path = "session_folder/presentation.pptx"
                    prs.save(presentation_path)
                    output_directory = 'session_folder'
                    os.system(f'libreoffice --headless --convert-to pdf --outdir {output_directory} {presentation_path}')
                    pdf_path = "session_folder/presentation.pdf"
                    audio_path = "session_folder/audio.mp3"
                    images = convert_from_path(pdf_path)
                    image_files = []

                    for i, img in enumerate(images):
                        img_path = f"session_folder/slide_{i}.png"
                        img.save(img_path, "PNG")
                        image_files.append(img_path)

                    audio_clip = AudioFileClip(audio_path)
                    video_clip = ImageSequenceClip(image_files, durations=[6] * len(images))
                    video_clip = video_clip.set_audio(audio_clip)
                    video_clip.fps = 24
                    output_path = f"session_folder/vid{k + 5 * count}.mp4"
                    video_clip.write_videofile(output_path, codec='libx264', audio_codec='aac')

                    for img_file in image_files:
                        os.remove(img_file)


                count += 1
    
        st.success("Combining all the videos into a single file... Please wait.")

        video_clips = []
        for i in range(0, 100):
            try:
                video_path = f"session_folder/vid{i}.mp4"
                video_clips.append(VideoFileClip(video_path))
            except:
                continue

        final_clip = concatenate_videoclips(video_clips)
        final_clip.write_videofile("session_folder/final.mp4")

        st.title("Final Video Player")
        final_video_path = "session_folder/final.mp4"

        with open(final_video_path, "rb") as f:
            video_bytes = f.read()
        st.video(video_bytes, format="video/mp4")

    st.header('Audio Speed')
    st.text('Audio at 1x speed (default)')
    audio_file = open('quality_1.mp3', 'rb')
    audio_bytes = audio_file.read()
    st.audio(audio_bytes, format='audio/mp3', start_time=0)
    st.text('Audio at 1.25x speed')
    audio_file = open('quality_1.25.mp3', 'rb')
    audio_bytes = audio_file.read()
    st.audio(audio_bytes, format='audio/mp3', start_time=0)
    st.text('Audio at 0.75x speed')
    audio_file = open('quality_0.75.mp3', 'rb')
    audio_bytes = audio_file.read()
    st.audio(audio_bytes, format='audio/mp3', start_time=0)
    st.header("voices")
    st.text("Alloy (default)")
    audio_file = open('quality_1.mp3', 'rb')
    audio_bytes = audio_file.read()
    st.audio(audio_bytes, format='audio/mp3', start_time=0)
    st.text("Echo")
    audio_file = open('quality_1_echo.mp3', 'rb')
    audio_bytes = audio_file.read()
    st.audio(audio_bytes, format='audio/mp3', start_time=0)
    st.text("Fable")
    audio_file = open('quality_1_fable.mp3', 'rb')
    audio_bytes = audio_file.read()
    st.audio(audio_bytes, format='audio/mp3', start_time=0)
    st.text("Onyx")
    audio_file = open('quality_1_onyx.mp3', 'rb')
    audio_bytes = audio_file.read()
    st.audio(audio_bytes, format='audio/mp3', start_time=0)
    st.text("Nova")
    audio_file = open('quality_1_nova.mp3', 'rb')
    audio_bytes = audio_file.read()
    st.audio(audio_bytes, format='audio/mp3', start_time=0)
    st.text("Shimmer")
    audio_file = open('quality_1_shimmer.mp3', 'rb')
    audio_bytes = audio_file.read()
    st.audio(audio_bytes, format='audio/mp3', start_time=0)

if __name__ == '__main__':
    app()
