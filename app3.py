import json
import os
import shutil
import streamlit as st
from openai import OpenAI
import requests
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import streamlit as st
from openai import OpenAI
from pptx.dml.color import RGBColor
from moviepy.editor import concatenate_videoclips, VideoFileClip, ImageSequenceClip, AudioFileClip
from pdf2image import convert_from_path

if 'openai_subheadings' not in st.session_state:
    st.session_state.openai_subheadings = ""


client = OpenAI(api_key="")

# functions
def gpt(system_prompt, user_prompt, model = 'gpt-3.5-turbo'):
    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        max_tokens = 2000
    )
    return response.choices[0].message.content

def get_text(lesson_url):
    response = requests.get(lesson_url)
    with open("session.pdf", 'wb') as f:
        f.write(response.content)
    loader = PyPDFLoader("session.pdf")
    documents = loader.load()
    text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=0)
    docs = text_splitter.split_documents(documents)
    text = ''.join([docs[i].page_content.replace("\t", " ") for i in range(len(docs))])
    with open("session.txt", 'w') as f:
        f.write(text)
    return text



# UI
if os.path.exists("session_folder"):
    shutil.rmtree("session_folder")
    os.mkdir("session_folder")
else:
    os.mkdir("session_folder")
system_message = """
    You are a teacher with decades of experience in teaching students of all levels and languages. You are 
    talented and help students to learn and understand the subject matter. You will be given a 
    lesson name, the book name, the age group of the student. You are tasked with creating a PPT for the students based on the lesson content you receive.
    Each slide of the ppt must explain a portion of the lesson perfectly. 
    Follow these guidelines strictly:
    Short description must be the overview of what you're about to teach and it has to be atleast 5 - 10 lines long.
    Long description must be the detailed explanation of the topic you're teaching and it has to be atleast 30 - 50 lines long.
    You will be given a list of sub headings. You must definetly explain these sub headings in your ppt according to the lesson.

    make sure to use vocabulary and language that is simple and very easy to understand.
    At the end of last slide, do not give any conclusion or summary. Just end the presentation. Thanks.
    Example: Slide <number>: <title> (this can be the topic you have chosen to teach.)
    short description: <short overview about the topic you're teaching, including the important keywords. This description has to be 5 - 10 lines"
    long description: <long description including the important keywords. This description has to be 30 - 50 lines or more. Explain precisely the topic you're explaining"
    
    DON'T LEAVE A LINE ANY WHERE IN YOUR RESPONSE.
    This is an example of how you should structure your response. You can use this as a template to structure your response.
[
    {
        "slide_title": "<slide title>",
        "short_description": "<short description>",
        "long_description": "<long description>"
    },

    {
        "slide_title": "<slide title>",
        "short_description": "<short description>",
        "long_description": "<long description>"
    },

    {
        "slide_title": "<slide title>",
        "short_description": "<short description>",
        "long_description": "<long description>"
    }

    
]
"""    

st.title("EKAI")
lesson_name = st.text_input("Enter the name of the lesson")
lesson_url = st.text_input("Enter the URL of the lesson")
grade = st.selectbox("Select the grade", ['Select Grade', "Grade 1", "Grade 2", "Grade 3", "Grade 4", "Grade 5", "Grade 6", "Grade 7", "Grade 8", "Grade 9", "Grade 10", "Grade 11", "Grade 12"])

if st.button("Generate subheadings"):
    with st.spinner("Generating..."):
        system_content = f"""
        You are an AI assistant specialized in preparing teaching material for students of all ages and grades. You will be given an entire lesson and you have to figure out the
        important subheadings and topics that will be useful for preparing material for the students. The subheadings and subtopics must be only from the lesson content you're provided. Don't make up your own.
        lesson name: {lesson_name}
        book name: NCERT Grade {grade}
        give them in a number format like this and not in any other format: 
        - <subheading>
        - <subheading>
        - <subheading>
        """
        user_content = f""" This is the lesson content: {get_text(lesson_url)} """
        subheadings = gpt(system_content, user_content, 'gpt-4o') #subheadings
        if 'openai_subheadings' in st.session_state and st.session_state.openai_subheadings == "":
            st.session_state.openai_subheadings = subheadings

openai_subheadings = st.session_state.openai_subheadings
user_subheadings = st.text_area("Edit your Subheadings", openai_subheadings.replace("  - ","- "), height=300)
if user_subheadings != openai_subheadings:
    st.session_state.openai_subheadings = user_subheadings

special_instructions = st.text_area("Special Instructions", height=100)
speed = st.select_slider('Select Speed', options=[0.75, 1, 1.25, 1.5], value=1)
voice = st.selectbox("Select a voice", ['Alloy', 'Echo', 'Fable', 'Onyx', 'Nova', 'Shimmer'])
count = 0
if st.button("Confirm"):
    user_subheadings = st.session_state.openai_subheadings
    final_subheadings = user_subheadings.replace("\n\n", "\n").replace("  - ","- ").replace("- ","").split("\n")
    chunks = [final_subheadings[i:i + 7] for i in range(0, len(final_subheadings), 7)]

    for chunk in chunks:
        subheadings = ""
        for i in chunk:
            subheadings += i + "\n"

        with open("session.txt", 'r') as f:
            script = f.read()

        user_message = f"""
        Lesson Name: {lesson_name}
        Book Name: "NCERT"
        grade: {grade}
        Sub headings: {subheadings}
        script: {script}
        You have to follow these special instructions: {special_instructions}
        """

        if len(user_message) / 4 > 512000:
            st.error("Please keep the script under 512000 characters")
        else:
            user_message = f"""
            Lesson Name: {lesson_name}
            book_name : NCERT
            Sub headings: {subheadings}
            Lesson content: {script}
        You have to follow these special instructions: {special_instructions}

            """

            st.success(f"Processing {count + 1} of {len(chunks)}")
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": f"{system_message}"},
                    {"role": "user", "content": f"{user_message}"},
                ],
                max_tokens=3000
            )

            full_script = response.choices[0].message.content
            full_script = full_script.replace("```","").replace("json", "")
            print("="*100)
            print(full_script)
            sections = eval(full_script)   
            for k in range(0, len(sections)):
                slide_title = sections[k]['slide_title']
                short_description = sections[k]['short_description']
                long_description = sections[k]['long_description']   
                speech_file_path = "session_folder/audio.mp3"
                try:
                    os.remove("session_folder/presentation.pdf")
                except:
                    pass

                try:
                    os.remove("session_folder/presentation.pptx")
                except:
                    pass

                try:
                    os.remove("audio.mp3")
                except:
                    print("failed to remove")

                with client.audio.speech.with_streaming_response.create(
                    model="tts-1",
                    voice=voice.lower(),
                    input=f"{long_description}",
                    speed=speed
                ) as response:
                    response.stream_to_file(speech_file_path)


                prs = Presentation()

                slide_background = RGBColor(0xD5, 0xE1, 0xDD)
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = slide_background

                title = slide.shapes.title
                title.text = str(slide_title)
                title.text_frame.paragraphs[0].font.bold = True  # Make the title bold
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # black color
                title.text_frame.paragraphs[0].font.size = Pt(28)  # Increase font size

                content = slide.placeholders[1]
                content.text = short_description
                content.text_frame.paragraphs[0].font.size = Pt(18)  # Set the font size to 18
                content.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x2E, 0x64, 0x4E)  # Forest green color

                fill = title.fill
                fill.gradient()
                fill.gradient_stops[0].color.rgb = RGBColor(0x9F, 0xDA, 0xC9)  # Light green
                fill.gradient_stops[1].color.rgb = RGBColor(0x2E, 0x64, 0x4E)  # Dark green

                content.shadow.inherit = False
                content.shadow.visible = True
                content.shadow.blur_radius = Pt(5)
                content.shadow.offset_x = Inches(0.1)
                content.shadow.offset_y = Inches(0.1)

                presentation_path = "session_folder/presentation.pptx"
                prs.save(presentation_path)
                output_directory = 'session_folder'
                os.system(f'libreoffice --headless --convert-to pdf --outdir {output_directory} {presentation_path}')
                pdf_path = "session_folder/presentation.pdf"
                audio_path = "session_folder/audio.mp3"
                images = convert_from_path(pdf_path)
                image_files = []

                for i, img in enumerate(images):
                    img_path = f"session_folder/slide_{i}.png"
                    img.save(img_path, "PNG")
                    image_files.append(img_path)

                audio_clip = AudioFileClip(audio_path)
                video_clip = ImageSequenceClip(image_files, durations=[6] * len(images))
                video_clip = video_clip.set_audio(audio_clip)
                video_clip.fps = 24
                output_path = f"session_folder/vid{k + 7 * count}.mp4"
                video_clip.write_videofile(output_path, codec='libx264', audio_codec='aac')

                for img_file in image_files:
                    os.remove(img_file)

            count += 1

    st.success("Done. We are now sending the videos to aws for combining all of them")
    st.balloons()
    
